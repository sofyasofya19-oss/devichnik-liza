#!/usr/bin/env python3
"""
Recreate the bachelorette party presentation in PowerPoint
with transitions, embedded videos, gradient backgrounds,
and entrance animations.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsmap
from PIL import Image, ImageDraw, ImageFilter
import os
import math

SLIDE_W = 13.333
SLIDE_H = 7.5
SLIDE_W_PX = 1920
SLIDE_H_PX = 1080

prs = Presentation()
prs.slide_width = Inches(SLIDE_W)
prs.slide_height = Inches(SLIDE_H)

PINK = RGBColor(0xf8, 0xb4, 0xd9)
VIOLET = RGBColor(0xa7, 0x8b, 0xfa)
GOLD = RGBColor(0xfb, 0xbf, 0x24)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_PURPLE = RGBColor(0xe9, 0xd5, 0xff)
DARK_PURPLE = RGBColor(0x1a, 0x0a, 0x2e)


def make_gradient(filename, colors, direction='diagonal'):
    """Create a gradient background image."""
    img = Image.new('RGB', (SLIDE_W_PX, SLIDE_H_PX))
    draw = ImageDraw.Draw(img)
    if direction == 'diagonal':
        for y in range(SLIDE_H_PX):
            for x in range(SLIDE_W_PX):
                t = (x / SLIDE_W_PX * 0.5
                     + y / SLIDE_H_PX * 0.5)
                r = int(colors[0][0] * (1 - t)
                        + colors[1][0] * t)
                g = int(colors[0][1] * (1 - t)
                        + colors[1][1] * t)
                b = int(colors[0][2] * (1 - t)
                        + colors[1][2] * t)
                draw.point((x, y), fill=(r, g, b))
    img.save(filename, 'PNG')
    return filename


def make_gradient_fast(filename, c1, c2, c3=None):
    """Fast 3-stop gradient via numpy-free approach."""
    img = Image.new('RGB', (SLIDE_W_PX, SLIDE_H_PX))
    pixels = img.load()
    for y in range(SLIDE_H_PX):
        for x in range(SLIDE_W_PX):
            t = (x + y) / (SLIDE_W_PX + SLIDE_H_PX)
            if c3 is None:
                r = int(c1[0] + (c2[0] - c1[0]) * t)
                g = int(c1[1] + (c2[1] - c1[1]) * t)
                b = int(c1[2] + (c2[2] - c1[2]) * t)
            else:
                if t < 0.5:
                    tt = t * 2
                    r = int(c1[0] + (c2[0] - c1[0]) * tt)
                    g = int(c1[1] + (c2[1] - c1[1]) * tt)
                    b = int(c1[2] + (c2[2] - c1[2]) * tt)
                else:
                    tt = (t - 0.5) * 2
                    r = int(c2[0] + (c3[0] - c2[0]) * tt)
                    g = int(c2[1] + (c3[1] - c2[1]) * tt)
                    b = int(c2[2] + (c3[2] - c2[2]) * tt)
            pixels[x, y] = (r, g, b)
    img.save(filename, 'PNG')
    return filename


# Pre-generate gradient backgrounds
os.makedirs('bg', exist_ok=True)
print("Generating gradient backgrounds...")
BG_MAIN = make_gradient_fast(
    'bg/main.png',
    (26, 10, 46), (107, 33, 168), (68, 49, 141)
)
BG_DARK = make_gradient_fast(
    'bg/dark.png',
    (15, 5, 26), (45, 20, 90), (26, 10, 46)
)
BG_TINDER = make_gradient_fast(
    'bg/tinder.png',
    (26, 10, 46), (107, 33, 168), (68, 49, 141)
)
BG_RESULT = make_gradient_fast(
    'bg/result.png',
    (21, 8, 37), (60, 20, 100), (30, 10, 50)
)
print("Backgrounds ready.")


def add_transition(slide, speed='med', trans_type='fade'):
    """Add a fade transition to a slide via XML injection."""
    xml = f'''
    <mc:AlternateContent
        xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006">
      <mc:Choice
        xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main"
        Requires="p14">
        <p:transition
          xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
          spd="{speed}" p14:dur="800" advClick="1">
          <p:{trans_type} />
        </p:transition>
      </mc:Choice>
      <mc:Fallback>
        <p:transition
          xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
          spd="{speed}" advClick="1">
          <p:{trans_type} />
        </p:transition>
      </mc:Fallback>
    </mc:AlternateContent>
    '''
    frag = parse_xml(xml)
    slide.element.insert(-1, frag)


def set_bg_image(slide, img_path):
    """Set slide background to an image."""
    slide.shapes.add_picture(
        img_path, Inches(0), Inches(0),
        width=prs.slide_width, height=prs.slide_height
    )


def add_text(slide, text, left, top, width, height,
             font_size=Pt(36), color=WHITE, bold=False,
             alignment=PP_ALIGN.CENTER, italic=False):
    """Add a textbox with styled text."""
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return tb


def add_title_slide(title, subtitle="", block_num="",
                    bg=None, transition='fade'):
    """Create a rich title/section slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_image(slide, bg or BG_MAIN)
    if block_num:
        add_text(slide, block_num, 0.5, 1.5, 12.3, 1,
                 Pt(26), VIOLET, bold=True)
    add_text(slide, title, 0.5, 2.8, 12.3, 2,
             Pt(54), WHITE, bold=True)
    if subtitle:
        add_text(slide, subtitle, 1, 4.8, 11.3, 2.2,
                 Pt(22), LIGHT_PURPLE)
    add_transition(slide, trans_type=transition)
    return slide


def add_question_slide(num, text):
    """Create a question slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_image(slide, BG_MAIN)
    add_text(slide, f"Вопрос {num} из 13",
             0.5, 1.2, 12.3, 0.8, Pt(22), VIOLET, True)
    add_text(slide, text, 1, 2.8, 11.3, 3.5,
             Pt(38), WHITE, True)
    add_text(slide, "Лиза, отвечай вслух!" if num == 1
             else "", 2, 6, 9, 0.8, Pt(18), PINK,
             italic=True)
    add_transition(slide)
    return slide


def add_video_slide(num):
    """Create a slide with embedded video."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_image(slide, BG_DARK)
    video_path = f"video/саша{num}.mp4"
    poster_path = f"posters/саша{num}.jpg"
    if os.path.exists(video_path) and \
            os.path.exists(poster_path):
        vid_w = Inches(8)
        vid_h = Inches(5)
        left = (prs.slide_width - vid_w) // 2
        top = Inches(0.8)
        slide.shapes.add_movie(
            video_path, left, top, vid_w, vid_h,
            poster_frame_image=poster_path,
            mime_type='video/mp4'
        )
    add_text(slide, f"Ответ Саши 💬",
             0.5, 6.3, 12.3, 0.8, Pt(20), PINK, True)
    add_transition(slide)
    return slide


def add_tinder_slide(name, fields, photo_path):
    """Create a tinder profile slide with photo."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_image(slide, BG_TINDER)
    if os.path.exists(photo_path):
        slide.shapes.add_picture(
            photo_path, Inches(0.6), Inches(1),
            width=Inches(4.2), height=Inches(5.2)
        )
    tb = slide.shapes.add_textbox(
        Inches(5.2), Inches(0.6),
        Inches(7.6), Inches(6.5)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = name
    run.font.size = Pt(30)
    run.font.color.rgb = PINK
    run.font.bold = True
    for label, value in fields:
        p = tf.add_paragraph()
        p.space_before = Pt(10)
        r = p.add_run()
        r.text = f"{label}: "
        r.font.size = Pt(17)
        r.font.color.rgb = VIOLET
        r.font.bold = True
        r2 = p.add_run()
        r2.text = value
        r2.font.size = Pt(17)
        r2.font.color.rgb = LIGHT_PURPLE
    add_transition(slide)
    return slide


def add_quote_slide(pair_num, quote_a, quote_b):
    """Create a quotes comparison slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_image(slide, BG_MAIN)
    add_text(
        slide,
        f"Пара {pair_num} — Какая цитата Сашина?",
        0.5, 0.4, 12.3, 0.8, Pt(24), VIOLET, True
    )
    add_text(slide, "Цитата А", 0.5, 1.5, 5.8, 0.6,
             Pt(18), GOLD, True)
    add_text(slide, quote_a, 0.5, 2.2, 5.8, 4.5,
             Pt(20), WHITE)
    add_text(slide, "Цитата Б", 6.8, 1.5, 5.8, 0.6,
             Pt(18), GOLD, True)
    add_text(slide, quote_b, 6.8, 2.2, 5.8, 4.5,
             Pt(20), WHITE)
    add_transition(slide)
    return slide


def add_result_slide(answer_text):
    """Create a result/answer slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_image(slide, BG_RESULT)
    add_text(slide, "🎉 Ответ:", 1, 1.5, 11.3, 1.5,
             Pt(40), GOLD, True)
    add_text(slide, answer_text, 1, 3.5, 11.3, 3.5,
             Pt(28), WHITE)
    add_transition(slide)
    return slide


def add_nhie_slide(text, bg_img=None):
    """Create a Never Have I Ever slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if bg_img and os.path.exists(bg_img):
        set_bg_image(slide, bg_img)
    else:
        set_bg_image(slide, BG_MAIN)
    add_text(slide, "Never have I ever...",
             1, 1.5, 11.3, 1.2, Pt(30), PINK, True)
    add_text(slide, text, 1, 3.5, 11.3, 3,
             Pt(34), WHITE, True)
    add_transition(slide)
    return slide


def add_hot_slide(photo_path):
    """Create a Hot or Not screenshot slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_image(slide, BG_TINDER)
    if os.path.exists(photo_path):
        img = Image.open(photo_path)
        w, h = img.size
        aspect = w / h
        max_h = Inches(6.2)
        max_w = Inches(6)
        if aspect > max_w / max_h:
            pic_w = max_w
            pic_h = int(max_w / aspect)
        else:
            pic_h = max_h
            pic_w = int(max_h * aspect)
        left = (prs.slide_width - pic_w) // 2
        top = Inches(0.5)
        slide.shapes.add_picture(
            photo_path, left, top,
            width=pic_w, height=pic_h
        )
    add_transition(slide)
    return slide


# ============================================
# BUILD PRESENTATION
# ============================================
print("Building slides...")

# --- TITLE SLIDE ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg_image(slide, BG_MAIN)
if os.path.exists('photos/title_1.png'):
    slide.shapes.add_picture(
        'photos/title_1.png', Inches(0), Inches(0),
        width=prs.slide_width, height=prs.slide_height
    )
add_text(slide, "✦ Девичник Лизы ✦", 0.5, 2.5, 12.3, 2,
         Pt(60), WHITE, True)
add_text(slide, "Она выходит замуж — и мы счастливы,\n"
         "что собрались по этому поводу.\n\nПоехали!",
         1, 4.8, 11.3, 2.2, Pt(24), LIGHT_PURPLE)
add_transition(slide, trans_type='fade')

# --- ROUND 1: ЛИЗА VS САША ---
add_title_slide(
    "Лиза vs Саша",
    "Лиза отвечает на вопрос —\n"
    "а потом смотрим, что думает Саша.",
    "Раунд 1"
)

questions = [
    "Кто из вас первый написал после "
    "первого свидания?",
    "Какое самое странное блюдо Лиза готовила?",
    "Из-за чего была ваша самая тупая ссора?",
    "Какую привычку Лизы ты сначала терпел, "
    "а потом полюбил?",
    "Что Лиза делает, когда думает, "
    "что ты не видишь?",
    "Кто из вас храпит?",
    "Какой фильм Лиза пересматривала "
    "больше всего?",
    "Что Черри любит больше — тебя или Лизу?",
    "Какой Лизин «контрол-фрик» момент "
    "тебя больше всего веселит?",
    "Какой подарок от Лизы был "
    "самым неожиданным?",
    "Если бы Лиза не стала продюсером — "
    "кем бы она была?",
    "Самый неловкий момент на свидании?",
    "Что бы ты сказал Лизе-из-прошлого, "
    "которая ещё тебя не знает?",
]

for i, q in enumerate(questions, 1):
    add_question_slide(i, q)
    add_video_slide(i)

# --- ROUND 2: ЛИЗА ВЫБИРАЕТ МУЖА ---
add_title_slide(
    "Лиза выбирает мужа",
    "Саша, конечно, очень хорош. Но мы нашли самых\n"
    "красивых мужчин в стране — ты просто обязана\n"
    "взглянуть.\n\nНадеемся, они заслужили свайп "
    "вправо напоследок от незамужней женщины",
    "Раунд 2"
)

candidates = [
    ("Кандидат №1 — Дмитрий, 30",
     "photos/dmitriy.png", [
        ("Город", "Москва (удалёнка, так что "
         "технически — диван)"),
        ("Работа", 'Senior-разработчик / '
         '"работаю 4 часа в день, остальное — думаю"'),
        ("О себе", '"Код пишу в халате. Дурь иногда '
         "помогает сконцентрироваться — это не я "
         "сказал, это Стив Джобс. Однажды починил "
         'прод в 3 ночи, не вставая с дивана."'),
        ("Плюс", "Зарплата 400к, философски "
         "относится к проблемам"),
        ("Минус", '"Щас доделаю" может означать '
         "от 5 минут до 3 дней"),
        ("Ищу", "Ту, которая не трогает мой второй "
         "монитор и не спрашивает зачем мне "
         "три клавиатуры"),
    ]),
    ("Кандидат №2 — Родриго, 28",
     "photos/rodrigo.png", [
        ("Город", "Мадрид"),
        ("Работа", "Таинственный предприниматель"),
        ("О себе", '"Люблю текилу, маленьких собак '
         'и женщин, которые всё контролируют"'),
        ("Плюс", "Сам готовит и убирает"),
        ("Минус", "Засыпает в 22:00"),
        ("Ищу", "Ту, которая составит мне "
         "расписание на жизнь"),
    ]),
    ("Кандидат №3 — Ахмед, 31",
     "photos/ahmed.png", [
        ("Город", "Нальчик → Москва"),
        ("Работа", '"Бизнес, сестра, не спрашивай"'),
        ("О себе", '"Увезу в горы. Там тихо, красиво,'
         ' и никто не найдёт. Шучу. Или нет."'),
        ("Плюс", "Решает любой вопрос за 10 минут"),
        ("Минус", "Мама всегда на первом месте"),
        ("Ищу", "Ту, за которую не стыдно "
         "перед родом"),
    ]),
    ("Кандидат №4 — Алехандро, 27",
     "photos/alejandro.png", [
        ("Город", "Барселона"),
        ("Работа", '"Художник... пишу, рисую, живу"'),
        ("О себе", '"Я покажу тебе такой закат, что '
         "ты забудешь как тебя зовут. Потом я тоже "
         'забуду, но это не важно."'),
        ("Плюс", "Говорит так, что хочется верить"),
        ("Минус", "Исчезает на 3 недели после "
         "идеального свидания"),
        ("Ищу", "Музу. На одну ночь. Максимум три."),
    ]),
    ("Кандидат №5 — Алижон, 38",
     "photos/alijon.png", [
        ("Город", "Ташкент → Москва (с 2009)"),
        ("Работа", 'Акушер-гинеколог / "делаю '
         'кесарево за 500$, гарантия"'),
        ("О себе", '"Руки золотые. 3000 детей '
         "принял. Чипсы с беконом не ем — харам. "
         'Но тебя на руках носить буду."'),
        ("Плюс", "Никогда не упадёт в обморок "
         "при виде крови"),
        ("Минус", "Все разговоры сводит к родам"),
        ("Ищу", "Ту, которая родит "
         "минимум четверых"),
    ]),
    ("Кандидат №6 — Геннадий, 45",
     "photos/gennadiy.png", [
        ("Город", "Воронеж (переезжать не планирует)"),
        ("Работа", 'Начальник отдела в "Водоканале"'),
        ("О себе", '"Могу достать любые трубы. '
         "Дача — 15 соток, баня с бассейном. "
         "По пятницам — шашлык. Люблю «Русский "
         'роман» и сериал «Гадалка»."'),
        ("Плюс", "Всё починит, всё построит, "
         "кран не течёт"),
        ("Минус", "В отпуск только в Анапу. "
         "Принципиально."),
        ("Ищу", "Хозяйственную. Чтоб борщ "
         "и чтоб не пилила."),
    ]),
]

for name, photo, fields in candidates:
    add_tinder_slide(name, fields, photo)

# --- ROUND 3: ЦИТАТЫ ---
add_title_slide(
    "Саша или Писатель?",
    "Две цитаты. Одна от Саши, другая от\n"
    "латиноамериканского писателя.\n"
    "Лиза угадывает. Ошиблась — шот 🥃",
    "Раунд 3"
)

quotes = [
    (
        "«Любовь — это не то, что находишь. "
        "Это то, что строишь каждый день заново.»",
        "«Нууууу, любовь — это, понимаете... типо "
        "мехового треугольника. Вот вечером лежим "
        "вместе и лежим спокойненько.»",
        "Цитата Б — это Саша!\n\n"
        "А цитата А — Габриэль Гарсиа Маркес"
    ),
    (
        "«Кто-то вот любит окрошку с квасом, "
        "а кто-то с кефиром. Чисто гипотетически. "
        "Но вы вместе всё равно, в одну сторону "
        "смотрите!»",
        "«Любовь так коротка, а забвение "
        "так длинно.»",
        "Цитата А — это Саша! Да, с окрошкой.\n\n"
        "А цитата Б — Пабло Неруда"
    ),
    (
        "«Люблю тебя не за то, кто ты, а за то, "
        "кем я становлюсь рядом с тобой.»",
        "«Люблю тебя потому что самая лучшая "
        "женщина во вселенной!»",
        "Цитата Б — это Саша! Коротко и по делу.\n\n"
        "А цитата А — Габриэль Гарсиа Маркес"
    ),
    (
        "«Я пью потому что ты пьёшь. "
        "Это любовь же и есть.»",
        "«Я пью, чтобы окружающие "
        "становились интереснее.»",
        "Цитата А — это Саша!\n"
        "Романтик и алкоголик в одном флаконе.\n\n"
        "А цитата Б — Эрнест Хемингуэй"
    ),
]

for i, (qa, qb, ans) in enumerate(quotes, 1):
    add_quote_slide(i, qa, qb)
    add_result_slide(ans)

# --- ROUND 4: NEVER HAVE I EVER ---
add_title_slide(
    "Never Have I Ever",
    "Кто делал — пьёт! 🍹",
    "Раунд 4"
)

nhie_items = [
    ("...напивалась и писала бывшему",
     "stock/couple1.png"),
    ("...просыпалась в чужой квартире и не помнила"
     " как туда попала", "stock/s1.png"),
    ("...пьяная признавалась кому-то в любви",
     "stock/s2.png"),
    ("...флиртовала ради бесплатного коктейля "
     "(и получала его)", "stock/s3.png"),
    ("...целовалась с кем-то просто потому что "
     "было скучно", "stock/s4.png"),
    ("...занималась сексом в общественном месте",
     "stock/s5.png"),
    ("...отправляла ню не тому человеку",
     "stock/s6.png"),
    ("...сталкерила человека настолько глубоко, "
     "что случайно лайкнула фото "
     "трёхлетней давности", "stock/s7.png"),
    ("...придумывала дебильную отмазку, чтобы "
     "уйти с плохого свидания", "stock/s8.png"),
    ("...блевала в такси и делала вид "
     "что всё нормально", "stock/s9.png"),
    ("...уходила из бара с незнакомцем, а подругам"
     " писала «я дома»", "stock/s10.png"),
    ("...плакала пьяная в туалете клуба, а через "
     "10 минут танцевала как ни в чём не бывало",
     "stock/s11.png"),
    ("...секстилась с человеком, с которым даже "
     "не встречалась вживую", "stock/couple3.png"),
]

for text, bg in nhie_items:
    add_nhie_slide(text, bg)

# --- FINAL SLIDE ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
if os.path.exists('photos/liza_sasha_1.png'):
    set_bg_image(slide, 'photos/liza_sasha_1.png')
else:
    set_bg_image(slide, BG_MAIN)
add_text(slide, "За Лизу и Сашу!", 0.5, 2.5, 12.3, 2,
         Pt(54), WHITE, True)
add_text(slide, "За любовь, за окрошку\n"
         "и за меховые треугольники! 🤍",
         1, 5, 11.3, 1.5, Pt(26), LIGHT_PURPLE)
add_transition(slide, trans_type='fade')

# --- P.S. HOT OR NOT ---
add_title_slide(
    "Hot or Not",
    "Бонус-раунд! Лиза, последний шанс оценить,\n"
    "от чего ты отказываешься.\n"
    "Реальные анкеты, реальный ужас. Свайпаем! 🔥",
    "P.S."
)

for i in range(1, 10):
    add_hot_slide(f"photos/hot_{i}.png")

# --- SAVE ---
output = "devichnik_liza.pptx"
prs.save(output)
print(f"\nDone! Saved: {output}")
print(f"Total slides: {len(prs.slides)}")
